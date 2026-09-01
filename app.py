"""PPTX -> Markdown 변환기 (오탈자 검토용)."""

import os
import sys
import queue
import subprocess
import threading
from pathlib import Path
from tkinter import filedialog

import customtkinter as ctk
from tkinterdnd2 import DND_FILES, TkinterDnD
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ctk.set_appearance_mode("System")
ctk.set_default_color_theme("blue")

APP_TITLE = "PPTX → Markdown 변환기"
OUTPUT_SUFFIX = "_오탈자검토"


def iter_shapes(shapes):
    """그룹 도형 내부까지 재귀적으로 순회한다."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def extract_slide_title(slide):
    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.has_text_frame:
        text = title_shape.text_frame.text.strip()
        if text:
            return text
    return ""


def shape_text_lines(shape):
    lines = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    return lines


def table_to_markdown(table):
    rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
    rows = [row for row in rows if any(cell for cell in row)]
    if not rows:
        return []
    md_lines = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * len(rows[0])) + " |"]
    for row in rows[1:]:
        md_lines.append("| " + " | ".join(row) + " |")
    return md_lines


def slide_notes_text(slide):
    if slide.has_notes_slide:
        notes_tf = slide.notes_slide.notes_text_frame
        if notes_tf is not None:
            return notes_tf.text.strip()
    return ""


def convert_slide(slide, index, log):
    title = extract_slide_title(slide)
    header = f"## Slide {index}: {title}" if title else f"## Slide {index}"

    title_shape = slide.shapes.title
    title_id = title_shape.shape_id if title_shape is not None else None

    shapes_sorted = sorted(
        iter_shapes(slide.shapes),
        key=lambda s: (getattr(s, "top", 0) or 0, getattr(s, "left", 0) or 0),
    )

    body_lines = []
    for shape in shapes_sorted:
        if title_id is not None and getattr(shape, "shape_id", None) == title_id:
            continue
        try:
            if shape.has_table:
                table_lines = table_to_markdown(shape.table)
                if table_lines:
                    body_lines.extend(table_lines)
                    body_lines.append("")
                continue
            if shape.has_text_frame:
                for line in shape_text_lines(shape):
                    body_lines.append(f"- {line}")
        except Exception as exc:
            log(f"  ! Slide {index} 도형 처리 중 오류 무시: {exc}")

    md = [header, ""]
    md.extend(body_lines if body_lines else ["_(텍스트 없음)_"])
    md.append("")

    notes = slide_notes_text(slide)
    if notes:
        md.append("> [발표자 노트]")
        md.extend(f"> {line.strip()}" for line in notes.splitlines() if line.strip())
        md.append("")

    return "\n".join(md)


def convert_pptx_to_markdown(pptx_path, log, progress_cb):
    prs = Presentation(pptx_path)
    total = len(prs.slides)
    lines = [f"# {Path(pptx_path).stem}", ""]
    for i, slide in enumerate(prs.slides, start=1):
        log(f"[{i}/{total}] 슬라이드 처리 중...")
        lines.append(convert_slide(slide, i, log))
        lines.append("---")
        lines.append("")
        progress_cb(i, total)
    return "\n".join(lines).rstrip() + "\n"


def build_output_path(pptx_path):
    p = Path(pptx_path)
    return p.with_name(f"{p.stem}{OUTPUT_SUFFIX}.md")


class App(ctk.CTk, TkinterDnD.DnDWrapper):
    def __init__(self):
        super().__init__()
        self.TkdndVersion = TkinterDnD._require(self)

        self.title(APP_TITLE)
        self.geometry("640x600")
        self.minsize(560, 500)

        self.selected_path = None
        self.output_path = None
        self.worker = None
        self.msg_queue = queue.Queue()

        self._build_ui()
        self.after(100, self._poll_queue)

    def _build_ui(self):
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(6, weight=1)

        ctk.CTkLabel(self, text=APP_TITLE, font=ctk.CTkFont(size=22, weight="bold")).grid(
            row=0, column=0, padx=24, pady=(24, 4), sticky="w"
        )
        ctk.CTkLabel(
            self,
            text="PPTX 파일을 깔끔한 Markdown으로 변환해 오탈자 검토를 쉽게 만들어드려요.",
            text_color=("gray30", "gray70"),
        ).grid(row=1, column=0, padx=24, pady=(0, 16), sticky="w")

        self.drop_frame = ctk.CTkFrame(
            self, height=140, corner_radius=16, border_width=2, border_color=("gray70", "gray40")
        )
        self.drop_frame.grid(row=2, column=0, padx=24, pady=8, sticky="ew")
        self.drop_frame.grid_propagate(False)
        self.drop_frame.grid_columnconfigure(0, weight=1)
        self.drop_frame.grid_rowconfigure(0, weight=1)

        self.drop_label = ctk.CTkLabel(
            self.drop_frame,
            text="여기로 PPTX 파일을 드래그하거나\n아래 버튼으로 선택하세요",
            font=ctk.CTkFont(size=14),
            justify="center",
        )
        self.drop_label.grid(row=0, column=0, sticky="nsew")

        for widget in (self.drop_frame, self.drop_label):
            widget.drop_target_register(DND_FILES)
            widget.dnd_bind("<<Drop>>", self._on_drop)

        button_row = ctk.CTkFrame(self, fg_color="transparent")
        button_row.grid(row=3, column=0, padx=24, pady=8, sticky="ew")
        button_row.grid_columnconfigure((0, 1), weight=1)

        self.select_btn = ctk.CTkButton(button_row, text="파일 선택", command=self._on_select_file)
        self.select_btn.grid(row=0, column=0, padx=(0, 6), sticky="ew")

        self.convert_btn = ctk.CTkButton(
            button_row, text="변환 시작", command=self._on_convert, state="disabled"
        )
        self.convert_btn.grid(row=0, column=1, padx=(6, 0), sticky="ew")

        self.file_label = ctk.CTkLabel(self, text="선택된 파일 없음", text_color=("gray40", "gray60"))
        self.file_label.grid(row=4, column=0, padx=24, pady=(4, 12), sticky="w")

        self.progress = ctk.CTkProgressBar(self)
        self.progress.set(0)
        self.progress.grid(row=5, column=0, padx=24, pady=(0, 12), sticky="ew")

        self.log_box = ctk.CTkTextbox(self, corner_radius=12)
        self.log_box.grid(row=6, column=0, padx=24, pady=(0, 12), sticky="nsew")
        self.log_box.configure(state="disabled")

        self.open_folder_btn = ctk.CTkButton(
            self, text="저장 폴더 열기", command=self._on_open_folder, state="disabled"
        )
        self.open_folder_btn.grid(row=7, column=0, padx=24, pady=(0, 24), sticky="ew")

    def _on_drop(self, event):
        paths = self.tk.splitlist(event.data)
        if not paths:
            return
        path = paths[0]
        if not path.lower().endswith(".pptx"):
            self._log(f"지원하지 않는 파일 형식입니다: {path}")
            return
        self._set_selected_file(path)

    def _on_select_file(self):
        path = filedialog.askopenfilename(
            title="PPTX 파일 선택", filetypes=[("PowerPoint 파일", "*.pptx")]
        )
        if path:
            self._set_selected_file(path)

    def _set_selected_file(self, path):
        self.selected_path = path
        self.output_path = None
        self.file_label.configure(text=path)
        self.convert_btn.configure(state="normal")
        self.open_folder_btn.configure(state="disabled")
        self.progress.set(0)
        self._clear_log()
        self._log(f"파일 선택됨: {path}")

    def _on_convert(self):
        if not self.selected_path or (self.worker and self.worker.is_alive()):
            return
        self.convert_btn.configure(state="disabled")
        self.select_btn.configure(state="disabled")
        self.open_folder_btn.configure(state="disabled")
        self.progress.set(0)
        self._clear_log()
        self.worker = threading.Thread(target=self._run_conversion, args=(self.selected_path,), daemon=True)
        self.worker.start()

    def _run_conversion(self, path):
        def log(msg):
            self.msg_queue.put(("log", msg))

        def progress_cb(done, total):
            self.msg_queue.put(("progress", done / total if total else 1))

        try:
            log(f"'{Path(path).name}' 변환을 시작합니다...")
            markdown = convert_pptx_to_markdown(path, log, progress_cb)
            output_path = build_output_path(path)
            output_path.write_text(markdown, encoding="utf-8")
            self.msg_queue.put(("done", str(output_path)))
        except Exception as exc:
            self.msg_queue.put(("error", str(exc)))

    def _poll_queue(self):
        try:
            while True:
                kind, payload = self.msg_queue.get_nowait()
                if kind == "log":
                    self._log(payload)
                elif kind == "progress":
                    self.progress.set(payload)
                elif kind == "done":
                    self.output_path = payload
                    self._log(f"완료! 저장 위치: {payload}")
                    self.progress.set(1)
                    self.convert_btn.configure(state="normal")
                    self.select_btn.configure(state="normal")
                    self.open_folder_btn.configure(state="normal")
                elif kind == "error":
                    self._log(f"오류 발생: {payload}")
                    self.convert_btn.configure(state="normal")
                    self.select_btn.configure(state="normal")
        except queue.Empty:
            pass
        self.after(100, self._poll_queue)

    def _log(self, message):
        self.log_box.configure(state="normal")
        self.log_box.insert("end", message + "\n")
        self.log_box.see("end")
        self.log_box.configure(state="disabled")

    def _clear_log(self):
        self.log_box.configure(state="normal")
        self.log_box.delete("1.0", "end")
        self.log_box.configure(state="disabled")

    def _on_open_folder(self):
        if not self.output_path:
            return
        folder = str(Path(self.output_path).parent)
        if sys.platform.startswith("win"):
            os.startfile(folder)
        elif sys.platform == "darwin":
            subprocess.Popen(["open", folder])
        else:
            subprocess.Popen(["xdg-open", folder])


if __name__ == "__main__":
    app = App()
    app.mainloop()
